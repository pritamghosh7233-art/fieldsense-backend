"""
combined_report_service.py
──────────────────────────
Wraps the two generator pipelines (PDF analyzer + PPT) for use by the
analyzed_reports router. Strips the custom-template logic from the PPT
generator and uses the random built-in themes only.
"""

import os
import json
import re
import random
import tempfile
import logging
import io
from pathlib import Path
from datetime import datetime

import boto3

logger = logging.getLogger(__name__)

BEDROCK_MODEL_ID = os.getenv("BEDROCK_MODEL_ID", "us.anthropic.claude-sonnet-4-6")
AWS_REGION       = os.getenv("AWS_REGION", "us-east-1")

# ─────────────────────────────────────────────────────────────────────────────
#  Bedrock client
# ─────────────────────────────────────────────────────────────────────────────

def _make_bedrock():
    return boto3.client(
        service_name="bedrock-runtime",
        region_name=AWS_REGION,
        aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
        aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY"),
    )


def _invoke(bedrock, system: str, user: str, max_tokens: int = 3000) -> str:
    resp = bedrock.invoke_model(
        modelId=BEDROCK_MODEL_ID,
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": max_tokens,
            "system": system,
            "messages": [{"role": "user", "content": user}],
        }),
    )
    return json.loads(resp["body"].read())["content"][0]["text"]


def _clean_json(raw: str) -> dict:
    raw = re.sub(r"^```json\s*", "", raw.strip(), flags=re.MULTILINE)
    raw = re.sub(r"^```\s*",     "", raw.strip(), flags=re.MULTILINE)
    raw = re.sub(r"```\s*$",     "", raw.strip(), flags=re.MULTILINE)
    raw = raw.strip()
    try:
        return json.loads(raw)
    except Exception:
        pass
    cleaned = re.sub(r',\s*([}\]])', r'\1', raw)
    try:
        return json.loads(cleaned)
    except Exception:
        pass
    match = re.search(r'\{.*\}', raw, re.DOTALL)
    if match:
        try:
            return json.loads(re.sub(r',\s*([}\]])', r'\1', match.group()))
        except Exception:
            pass
    raise ValueError(f"Cannot parse JSON: {raw[:300]}")


# ─────────────────────────────────────────────────────────────────────────────
#  PDF pipeline  (uses the analyzer logic from the provided analyzer_v2 code)
# ─────────────────────────────────────────────────────────────────────────────

ANALYZER_SYSTEM = """You are an elite industrial inspection analyst. You are given combined
approved inspection reports from multiple sessions across various zones and time periods.
Synthesise them into a single structured analysis JSON for a professional PDF report.

Return ONLY valid JSON — no markdown fences, no prose outside the JSON.

Schema:
{
  "report_title": "string",
  "cover_heading": "string — 2-3 sentences summarising scope",
  "document_type": "Combined Inspection Analysis Report",
  "domain": "string — e.g. 'Industrial / Manufacturing'",
  "prepared_for": "string",
  "key_metrics": [
    {"label":"string","value":"string","change":"string","status":"GOOD|WARNING|CRITICAL"}
  ],
  "snapshot_bullets": ["string — must include a number or specific finding"],
  "key_findings": [
    {"category":"string","bullets":["string"]}
  ],
  "data_tables": [
    {"title":"string","headers":["string"],"rows":[["string"]],"highlight_rows":[0]}
  ],
  "charts_data": [
    {"chart_type":"line|bar|stacked_bar|pie|area","title":"string","x_label":"string",
     "y_label":"string","series":[{"name":"string","data":[number]}],"labels":["string"]}
  ],
  "risk_matrix": [
    {"risk":"string","likelihood":"HIGH|MEDIUM|LOW","impact":"HIGH|MEDIUM|LOW","mitigation":"string"}
  ],
  "recommendations": [
    {"priority":"IMMEDIATE|SHORT-TERM|MEDIUM-TERM","action":"string","bullets":["string"],"timeline":"string"}
  ],
  "projections": [],
  "conclusion": {
    "query_addressed":"string",
    "direct_answer":"string",
    "bullets":[{"point":"string","evidence":"string"}],
    "paragraph":"string"
  }
}"""


def _analyze_with_bedrock(combined_text: str, user_query: str = "") -> dict:
    bedrock = _make_bedrock()
    prompt = f"USER QUERY: {user_query}\n\nINSPECTION REPORTS:\n{combined_text}"
    raw = _invoke(bedrock, ANALYZER_SYSTEM, prompt, max_tokens=6000)
    return _clean_json(raw)


def run_pdf_pipeline(combined_text: str, user_query: str, output_path: str):
    """Generate a combined analyzed PDF report from inspection text."""
    logger.info("[pdf_pipeline] Analyzing combined text with Bedrock...")
    analysis = _analyze_with_bedrock(combined_text, user_query)

    logger.info("[pdf_pipeline] Generating charts...")
    charts = _generate_charts(analysis)

    logger.info("[pdf_pipeline] Building PDF...")
    _build_pdf(analysis, charts, output_path)

    # Cleanup chart temp files
    for p in charts.values():
        try:
            os.unlink(p)
        except Exception:
            pass
    logger.info(f"[pdf_pipeline] Done → {output_path}")


# ─────────────────────────────────────────────────────────────────────────────
#  Chart generation (matplotlib)
# ─────────────────────────────────────────────────────────────────────────────

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

CHART_COLORS = ['#0288D1','#FF6F00','#2E7D32','#6A1B9A','#00838F','#AD1457','#558B2F','#4527A0']

P = {
    'primary': '#1A237E', 'accent': '#0288D1', 'accent_light': '#E3F2FD',
    'success': '#2E7D32', 'success_bg': '#E8F5E9', 'warning': '#E65100',
    'warning_bg': '#FFF3E0', 'critical': '#B71C1C', 'critical_bg': '#FFEBEE',
    'bg_light': '#F5F7FA', 'border': '#CFD8DC', 'text': '#1C1C1E', 'muted': '#607D8B',
}


def _save_fig(fig) -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return tmp.name


def _style_ax(ax, title, x_label="", y_label=""):
    ax.set_title(title, fontsize=11, fontweight='bold', pad=10, color=P['text'])
    if x_label: ax.set_xlabel(x_label, fontsize=8.5, color=P['muted'])
    if y_label: ax.set_ylabel(y_label, fontsize=8.5, color=P['muted'])
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.tick_params(colors=P['muted'], labelsize=7.5)
    ax.grid(True, alpha=0.18, linestyle='--', color=P['muted'])


def _make_bar_chart(c, stacked=False) -> str:
    fig, ax = plt.subplots(figsize=(7.5, 3.4), facecolor='white')
    labels = c.get("labels", [])
    series_list = c.get("series", [])
    n = max(len(labels) or 1, max((len(s.get("data",[])) for s in series_list), default=1))
    x = np.arange(n)
    width = 0.65 / max(len(series_list), 1) if not stacked else 0.6
    bottoms = np.zeros(n)
    for i, s in enumerate(series_list):
        col = CHART_COLORS[i % len(CHART_COLORS)]
        vals = np.array(s.get("data", [0]*n)[:n])
        if stacked:
            ax.bar(x, vals, width, bottom=bottoms, label=s.get("name"), color=col, alpha=0.85)
            bottoms += vals
        else:
            offset = (i - len(series_list)/2 + 0.5) * width
            ax.bar(x + offset, vals, width, label=s.get("name"), color=col, alpha=0.85)
    if labels:
        ax.set_xticks(x); ax.set_xticklabels(labels, rotation=15, ha='right', fontsize=7.5)
    if len(series_list) > 1:
        ax.legend(fontsize=7.5, framealpha=0.5)
    _style_ax(ax, c.get("title",""), c.get("x_label",""), c.get("y_label",""))
    plt.tight_layout(pad=0.4)
    return _save_fig(fig)


def _make_line_chart(c) -> str:
    fig, ax = plt.subplots(figsize=(7.5, 3.4), facecolor='white')
    labels = c.get("labels", [])
    for i, s in enumerate(c.get("series", [])):
        col = CHART_COLORS[i % len(CHART_COLORS)]
        vals = s.get("data", [])
        ax.plot(range(len(vals)), vals, marker='o', lw=2.2, ms=5.5, color=col, label=s.get("name"))
        ax.fill_between(range(len(vals)), vals, alpha=0.07, color=col)
    if labels:
        ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, rotation=15, ha='right', fontsize=7.5)
    ax.legend(fontsize=7.5, framealpha=0.5)
    _style_ax(ax, c.get("title",""), c.get("x_label",""), c.get("y_label",""))
    plt.tight_layout(pad=0.4)
    return _save_fig(fig)


def _make_pie_chart(c) -> str:
    fig, (ax, ax_leg) = plt.subplots(1, 2, figsize=(7.5, 3.4), facecolor='white',
                                      gridspec_kw={'width_ratios': [3, 1.2]})
    series = c.get("series", [{}])[0]
    vals = series.get("data", [])
    labels = c.get("labels", [f"Item {i+1}" for i in range(len(vals))])
    colors = CHART_COLORS[:len(vals)]
    wedges, _, autotexts = ax.pie(vals, autopct='%1.1f%%', colors=colors, startangle=140,
                                   pctdistance=0.78, wedgeprops={'linewidth':1.5,'edgecolor':'white'})
    for at in autotexts:
        at.set_fontsize(7.5); at.set_fontweight('bold'); at.set_color('white')
    ax.set_title(c.get("title",""), fontsize=11, fontweight='bold', pad=10)
    ax_leg.axis('off')
    total = sum(vals)
    for i, (lbl, val) in enumerate(zip(labels, vals)):
        pct = val/total*100 if total else 0
        y = 0.9 - i * 0.22
        ax_leg.text(0.04, y, '■', color=colors[i % len(colors)], fontsize=11, transform=ax_leg.transAxes, va='center')
        ax_leg.text(0.18, y, f'{lbl}', color=P['text'], fontsize=8, transform=ax_leg.transAxes, va='center', fontweight='bold')
        ax_leg.text(0.18, y-0.1, f'{pct:.1f}%', color=P['muted'], fontsize=7.5, transform=ax_leg.transAxes, va='center')
    plt.tight_layout(pad=0.4)
    return _save_fig(fig)


def _make_risk_matrix(risks) -> str:
    fig, ax = plt.subplots(figsize=(7.5, 4.8), facecolor='white')
    levels = {"LOW": 0, "MEDIUM": 1, "HIGH": 2}
    grid_colors = [['#C8E6C9','#FFF9C4','#FFCCBC'],['#FFF9C4','#FFCCBC','#FFCDD2'],['#FFCCBC','#FFCDD2','#EF9A9A']]
    for row in range(3):
        for col in range(3):
            ax.add_patch(plt.Rectangle([col, row], 1, 1, color=grid_colors[row][col], ec='white', lw=2))
    placed = {}
    for risk in risks:
        lk = levels.get(risk.get("likelihood","MEDIUM").upper(), 1)
        im = levels.get(risk.get("impact","MEDIUM").upper(), 1)
        key = (lk, im)
        offset = placed.get(key, 0) * 0.24
        placed[key] = placed.get(key, 0) + 1
        ax.annotate(risk.get("risk","")[:24], xy=(lk+0.5, im+0.62-offset),
                    ha='center', va='center', fontsize=7.5, fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.22', fc='white', alpha=0.93, ec='#90A4AE', lw=1))
    ax.set_xlim(0,3); ax.set_ylim(0,3)
    ax.set_xticks([0.5,1.5,2.5]); ax.set_xticklabels(['Low','Medium','High'], fontsize=9)
    ax.set_yticks([0.5,1.5,2.5]); ax.set_yticklabels(['Low','Medium','High'], fontsize=9)
    ax.set_xlabel("Likelihood  →", fontsize=10); ax.set_ylabel("Impact  →", fontsize=10)
    ax.set_title("Risk Assessment Matrix", fontsize=12, fontweight='bold', pad=10)
    patches = [mpatches.Patch(color='#C8E6C9',label='Low'),mpatches.Patch(color='#FFF9C4',label='Medium'),
               mpatches.Patch(color='#FFCCBC',label='High'),mpatches.Patch(color='#EF9A9A',label='Critical')]
    ax.legend(handles=patches, loc='upper left', bbox_to_anchor=(1.02,1), fontsize=7.5)
    plt.tight_layout(pad=0.4)
    return _save_fig(fig)


def _generate_charts(analysis: dict) -> dict:
    charts = {}
    for c in analysis.get("charts_data", []):
        ctype = c.get("chart_type","bar").lower()
        try:
            if ctype == "line": path = _make_line_chart(c)
            elif ctype == "pie": path = _make_pie_chart(c)
            elif ctype == "stacked_bar": path = _make_bar_chart(c, stacked=True)
            else: path = _make_bar_chart(c)
            charts[c.get("title", f"chart_{len(charts)}")] = path
        except Exception as e:
            logger.warning(f"Chart failed: {e}")
    if analysis.get("risk_matrix"):
        try:
            charts["__risk_matrix__"] = _make_risk_matrix(analysis["risk_matrix"])
        except Exception as e:
            logger.warning(f"Risk matrix failed: {e}")
    return charts


# ─────────────────────────────────────────────────────────────────────────────
#  PDF builder (ReportLab)
# ─────────────────────────────────────────────────────────────────────────────

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    Image as RLImage, HRFlowable, PageBreak, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.pdfgen import canvas as rl_canvas

PAGE_W, PAGE_H = A4
MARGIN = 1.8 * cm


def _ps(name, **kw): return ParagraphStyle(name, **kw)


class _NumberedCanvas(rl_canvas.Canvas):
    def __init__(self, *args, **kwargs):
        rl_canvas.Canvas.__init__(self, *args, **kwargs)
        self._saved_pages = []

    def showPage(self):
        self._saved_pages.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        for i, page in enumerate(self._saved_pages):
            self.__dict__.update(page)
            if i >= 1:
                self._draw_footer(i + 1)
            rl_canvas.Canvas.showPage(self)
        rl_canvas.Canvas.save(self)

    def _draw_footer(self, page_num):
        self.saveState()
        self.setFont('Helvetica', 7.5)
        self.setFillColorRGB(0.45, 0.45, 0.45)
        self.drawString(MARGIN, 1.0*cm, f"FieldSense · Combined Inspection Analysis · {datetime.now().strftime('%Y-%m-%d')}")
        self.drawRightString(PAGE_W - MARGIN, 1.0*cm, f"Page {page_num}")
        self.setStrokeColor(rl_colors.HexColor(P['accent']))
        self.setLineWidth(1)
        self.line(MARGIN, 1.35*cm, PAGE_W - MARGIN, 1.35*cm)
        self.restoreState()


def _build_pdf(analysis: dict, charts: dict, output_path: str):
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            leftMargin=MARGIN, rightMargin=MARGIN,
                            topMargin=1.8*cm, bottomMargin=2.0*cm)
    story = []

    S = {
        'h1': _ps('h1', fontSize=20, fontName='Helvetica-Bold',
                  textColor=rl_colors.HexColor(P['primary']), spaceAfter=6),
        'h2': _ps('h2', fontSize=13, fontName='Helvetica-Bold',
                  textColor=rl_colors.HexColor(P['accent']), spaceBefore=8, spaceAfter=6),
        'body': _ps('body', fontSize=9.5, fontName='Helvetica',
                    textColor=rl_colors.HexColor(P['text']), leading=14),
        'bullet': _ps('bullet', fontSize=9, fontName='Helvetica',
                      textColor=rl_colors.HexColor(P['text']), leftIndent=12, spaceAfter=2, leading=13),
        'th': _ps('th', fontSize=8.5, fontName='Helvetica-Bold',
                  textColor=rl_colors.white, alignment=TA_CENTER),
    }

    # ── Cover-ish header ─────────────────────────────────────────────────────
    story.append(Paragraph(analysis.get("report_title", "Combined Inspection Analysis"), S['h1']))
    story.append(HRFlowable(width="100%", thickness=2, color=rl_colors.HexColor(P['accent']), spaceAfter=8))
    story.append(Paragraph(analysis.get("cover_heading", ""), S['body']))
    story.append(Spacer(1, 12))

    # ── Key metrics ──────────────────────────────────────────────────────────
    metrics = analysis.get("key_metrics", [])
    if metrics:
        story.append(Paragraph("Key Metrics", S['h2']))
        status_cfg = {
            'GOOD': (P['success'], P['success_bg']),
            'WARNING': (P['warning'], P['warning_bg']),
            'CRITICAL': (P['critical'], P['critical_bg']),
        }
        chunk_size = 4
        for chunk in [metrics[i:i+chunk_size] for i in range(0, len(metrics), chunk_size)]:
            cells = []
            for m in chunk:
                tc, bg = status_cfg.get(m.get("status","GOOD").upper(), (P['text'], P['bg_light']))
                change = m.get("change","")
                arrow = "▲" if "+" in change else "▼" if "-" in change else "●"
                cells.append([
                    Paragraph(m.get("label","").upper(), _ps('ml', fontSize=7.5, fontName='Helvetica',
                               textColor=rl_colors.HexColor(P['muted']), spaceAfter=2)),
                    Paragraph(m.get("value",""), _ps('mv', fontSize=18, fontName='Helvetica-Bold',
                               textColor=rl_colors.HexColor(tc), leading=22, spaceAfter=2)),
                    Paragraph(f'{arrow} {change}', _ps('mc', fontSize=8, fontName='Helvetica-Bold',
                               textColor=rl_colors.HexColor(tc))),
                ])
            while len(cells) < chunk_size:
                cells.append([Paragraph("", S['body'])])
            col_w = doc.width / chunk_size
            t = Table([cells], colWidths=[col_w]*chunk_size)
            ts_list = [
                ('VALIGN',(0,0),(-1,-1),'TOP'), ('TOPPADDING',(0,0),(-1,-1),10),
                ('BOTTOMPADDING',(0,0),(-1,-1),10), ('LEFTPADDING',(0,0),(-1,-1),10),
                ('BOX',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
                ('INNERGRID',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
            ]
            for i, m in enumerate(chunk):
                tc, bg = status_cfg.get(m.get("status","GOOD").upper(), (P['text'], P['bg_light']))
                ts_list.append(('BACKGROUND',(i,0),(i,0),rl_colors.HexColor(bg)))
                ts_list.append(('LINEABOVE',(i,0),(i,0),3,rl_colors.HexColor(tc)))
            t.setStyle(TableStyle(ts_list))
            story.append(t); story.append(Spacer(1,6))

    # ── Snapshot bullets ─────────────────────────────────────────────────────
    bullets = analysis.get("snapshot_bullets",[])
    if bullets:
        story.append(Paragraph("At-a-Glance Snapshot", S['h2']))
        rows = [[Paragraph(f'▸  {b}', S['bullet'])] for b in bullets]
        t = Table(rows, colWidths=[doc.width])
        t.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,-1),rl_colors.HexColor(P['accent_light'])),
            ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
            ('LEFTPADDING',(0,0),(-1,-1),12),
            ('LINEBEFORE',(0,0),(-1,-1),3,rl_colors.HexColor(P['primary'])),
            ('BOX',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
        ]))
        story.append(t); story.append(Spacer(1,10))

    story.append(PageBreak())

    # ── Key findings ─────────────────────────────────────────────────────────
    findings = analysis.get("key_findings",[])
    if findings:
        story.append(Paragraph("Key Findings", S['h2']))
        cat_colors = [P['accent'],P['primary'],'#6A1B9A','#00838F']
        bg_colors  = [P['accent_light'],'#E8EAF6','#F3E5F5','#E0F7FA']
        for i, f in enumerate(findings):
            cc = cat_colors[i % len(cat_colors)]
            bc = bg_colors[i % len(bg_colors)]
            block = [Paragraph(f'<b>{f.get("category","")}</b>',
                               _ps('ch',fontSize=10,fontName='Helvetica-Bold',
                                   textColor=rl_colors.HexColor(cc)))]
            for b in f.get("bullets",[]):
                block.append(Paragraph(f'▪  {b}', S['bullet']))
            rows_t = [[item] for item in block]
            t = Table(rows_t, colWidths=[doc.width])
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),rl_colors.HexColor(bc)),
                ('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
                ('LEFTPADDING',(0,0),(-1,-1),12),
                ('LINEBEFORE',(0,0),(-1,-1),3.5,rl_colors.HexColor(cc)),
                ('BOX',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
            ]))
            story.append(KeepTogether([t, Spacer(1,6)]))

    # ── Charts ───────────────────────────────────────────────────────────────
    if charts:
        story.append(PageBreak())
        story.append(Paragraph("Data Visualizations", S['h2']))
        chart_data = analysis.get("charts_data",[])
        ordered = [(c.get("title",""), charts[c.get("title","")]) for c in chart_data if c.get("title","") in charts]
        if "__risk_matrix__" in charts:
            ordered.append(("Risk Assessment Matrix", charts["__risk_matrix__"]))
        col_w = doc.width / 2 - 3
        for i in range(0, len(ordered), 2):
            pair = ordered[i:i+2]
            row_cells = [[RLImage(path, width=col_w, height=2.45*inch)] for _, path in pair]
            while len(row_cells) < 2:
                row_cells.append([Paragraph("", S['body'])])
            t = Table([row_cells], colWidths=[col_w+3, col_w+3])
            t.setStyle(TableStyle([
                ('VALIGN',(0,0),(-1,-1),'TOP'),('TOPPADDING',(0,0),(-1,-1),7),
                ('BACKGROUND',(0,0),(-1,-1),rl_colors.HexColor(P['bg_light'])),
                ('BOX',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
            ]))
            story.append(t); story.append(Spacer(1,6))

    # ── Data tables ──────────────────────────────────────────────────────────
    tables = analysis.get("data_tables",[])
    if tables:
        story.append(PageBreak())
        story.append(Paragraph("Data Tables", S['h2']))
        for dt in tables:
            story.append(Paragraph(dt.get("title",""), _ps('dt', fontSize=10, fontName='Helvetica-Bold',
                                   textColor=rl_colors.HexColor(P['accent']), spaceAfter=4)))
            headers = dt.get("headers",[])
            if not headers: continue
            n_cols = len(headers)
            col_w = doc.width / n_cols
            data = [[Paragraph(h, S['th']) for h in headers]]
            for row in dt.get("rows",[]):
                data.append([Paragraph(str(v), _ps('tc', fontSize=8.5, fontName='Helvetica',
                                        textColor=rl_colors.HexColor(P['text']))) for v in row])
            t = Table(data, colWidths=[col_w]*n_cols, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,0),rl_colors.HexColor(P['primary'])),
                ('TEXTCOLOR',(0,0),(-1,0),rl_colors.white),
                ('GRID',(0,0),(-1,-1),0.4,rl_colors.HexColor(P['border'])),
                ('ROWBACKGROUNDS',(0,1),(-1,-1),[rl_colors.white, rl_colors.HexColor(P['bg_light'])]),
                ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                ('LEFTPADDING',(0,0),(-1,-1),6),
            ]))
            story.append(KeepTogether([t, Spacer(1,10)]))

    # ── Recommendations ──────────────────────────────────────────────────────
    recs = analysis.get("recommendations",[])
    if recs:
        story.append(Paragraph("Recommendations", S['h2']))
        pri_cfg = {
            'IMMEDIATE': (P['critical'], P['critical_bg']),
            'SHORT-TERM': (P['warning'], P['warning_bg']),
            'MEDIUM-TERM': (P['success'], P['success_bg']),
        }
        for r in recs:
            pri = r.get("priority","MEDIUM-TERM").upper()
            tc, bg = pri_cfg.get(pri, (P['text'], P['bg_light']))
            block = [Paragraph(f'<b>[{pri}]</b>  {r.get("action","")}  · {r.get("timeline","")}',
                               _ps('rh', fontSize=10, fontName='Helvetica', leading=14))]
            for b in r.get("bullets",[]):
                block.append(Paragraph(f'▪  {b}', S['bullet']))
            rows_t = [[item] for item in block]
            t = Table(rows_t, colWidths=[doc.width])
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),rl_colors.HexColor(bg)),
                ('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
                ('LEFTPADDING',(0,0),(-1,-1),12),
                ('LINEBEFORE',(0,0),(-1,-1),4,rl_colors.HexColor(tc)),
                ('BOX',(0,0),(-1,-1),0.5,rl_colors.HexColor(P['border'])),
            ]))
            story.append(KeepTogether([t, Spacer(1,6)]))

    # ── Conclusion ───────────────────────────────────────────────────────────
    conclusion = analysis.get("conclusion",{})
    if conclusion:
        story.append(PageBreak())
        story.append(Paragraph("Conclusion", S['h2']))
        if conclusion.get("direct_answer"):
            story.append(Paragraph(conclusion["direct_answer"], S['body']))
            story.append(Spacer(1,8))
        if conclusion.get("paragraph"):
            story.append(Paragraph(conclusion["paragraph"], S['body']))

    doc.build(story, canvasmaker=_NumberedCanvas)


# ─────────────────────────────────────────────────────────────────────────────
#  PPT pipeline  (random themes only — no custom template)
# ─────────────────────────────────────────────────────────────────────────────

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)

ALL_THEMES = {
    "corporate_navy": {
        "name": "Corporate Navy",
        "bg_dark": PptxRGB(0x1E,0x27,0x61), "bg_light": PptxRGB(0xF4,0xF6,0xFB),
        "accent": PptxRGB(0x00,0xA8,0xE8), "accent2": PptxRGB(0xF9,0x61,0x67),
        "text_light": PptxRGB(0xFF,0xFF,0xFF), "text_dark": PptxRGB(0x1A,0x1A,0x2E),
        "text_muted": PptxRGB(0x66,0x72,0x8A), "footer_bg": PptxRGB(0x1E,0x27,0x61),
        "title_font": "Calibri", "body_font": "Calibri Light",
        "chart_colors": ["#00A8E8","#F96167","#1E2761","#F9E795","#02C39A"],
    },
    "emerald_executive": {
        "name": "Emerald Executive",
        "bg_dark": PptxRGB(0x0D,0x3B,0x2B), "bg_light": PptxRGB(0xF0,0xF7,0xF4),
        "accent": PptxRGB(0x2E,0xCC,0x8B), "accent2": PptxRGB(0xF5,0xA6,0x23),
        "text_light": PptxRGB(0xFF,0xFF,0xFF), "text_dark": PptxRGB(0x0D,0x1F,0x17),
        "text_muted": PptxRGB(0x5A,0x7A,0x6A), "footer_bg": PptxRGB(0x0D,0x3B,0x2B),
        "title_font": "Georgia", "body_font": "Calibri Light",
        "chart_colors": ["#2ECC8B","#F5A623","#0D3B2B","#A8E6CF","#FF6B6B"],
    },
    "steel_amber": {
        "name": "Steel & Amber",
        "bg_dark": PptxRGB(0x2B,0x2D,0x42), "bg_light": PptxRGB(0xF5,0xF5,0xF0),
        "accent": PptxRGB(0xEF,0xC0,0x50), "accent2": PptxRGB(0xEF,0x82,0x50),
        "text_light": PptxRGB(0xFF,0xFF,0xFF), "text_dark": PptxRGB(0x2B,0x2D,0x42),
        "text_muted": PptxRGB(0x7A,0x7C,0x8A), "footer_bg": PptxRGB(0x2B,0x2D,0x42),
        "title_font": "Century Gothic", "body_font": "Century Gothic",
        "chart_colors": ["#EFC050","#EF8250","#2B2D42","#8D99AE","#D90429"],
    },
    "arctic_minimal": {
        "name": "Arctic Minimal",
        "bg_dark": PptxRGB(0x1C,0x3A,0x4A), "bg_light": PptxRGB(0xF7,0xFB,0xFF),
        "accent": PptxRGB(0x00,0xD4,0xD4), "accent2": PptxRGB(0xFF,0x8C,0x42),
        "text_light": PptxRGB(0xFF,0xFF,0xFF), "text_dark": PptxRGB(0x0D,0x1B,0x2A),
        "text_muted": PptxRGB(0x5C,0x7A,0x8A), "footer_bg": PptxRGB(0x1C,0x3A,0x4A),
        "title_font": "Arial", "body_font": "Arial Narrow",
        "chart_colors": ["#00D4D4","#FF8C42","#1C3A4A","#A8DADC","#E63946"],
    },
}

_ACTIVE_THEME = ALL_THEMES["corporate_navy"]

FOOTER_H = Inches(0.32)

PPT_SLIDE_SYSTEM = """You are an expert industrial presentation designer.
Given combined inspection report data, produce a concise professional slide plan.
Return ONLY valid JSON — no markdown fences.
Schema:
{
  "title": "string", "subtitle": "string", "author": "string",
  "slides": [
    {"type":"title","title":"string","subtitle":"string"},
    {"type":"toc","title":"Agenda","items":["string"]},
    {"type":"section","title":"string"},
    {"type":"content","title":"string","bullets":["string"]},
    {"type":"table","title":"string","headers":["string"],"rows":[["string"]]},
    {"type":"bar_chart","title":"string","chart_title":"string","x_label":"string","y_label":"string",
     "categories":["string"],"series":[{"name":"string","values":[number]}],"data_labels":true,"summary_bullets":["string"]},
    {"type":"conclusion","title":"Conclusion","bullets":["string"]},
    {"type":"thankyou","title":"Thank You","message":"string"}
  ]
}
Rules: 12-18 slides total. Only use real numbers from the data. No fabricated chart values.
"""


def _ppt_invoke(combined_text: str, user_query: str) -> dict:
    bedrock = _make_bedrock()
    prompt = f"USER QUERY: {user_query}\n\nINSPECTION DATA:\n{combined_text[:50000]}"
    raw = _invoke(bedrock, PPT_SLIDE_SYSTEM, prompt, max_tokens=5000)
    return _clean_json(raw)


def _new_blank_slide(prs, bg_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide


def _add_textbox(slide, text, x, y, w, h, font_name, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_shape(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_footer(slide, slide_num):
    T = _ACTIVE_THEME
    footer_y = T.get("slide_h", _SLIDE_H) - FOOTER_H
    _add_shape(slide, Inches(0), footer_y, _SLIDE_W, FOOTER_H, T["footer_bg"])
    _add_textbox(slide, f"FieldSense · Slide {slide_num}", Inches(0.3),
                 footer_y + Inches(0.05), Inches(4), FOOTER_H - Inches(0.06),
                 T["body_font"], 9, color=T["text_muted"])


def _add_header(slide, title):
    T = _ACTIVE_THEME
    _add_shape(slide, Inches(0), Inches(0), _SLIDE_W, Inches(1.3), T["bg_dark"])
    _add_textbox(slide, title, Inches(0.6), Inches(0.22), Inches(11.5), Inches(0.85),
                 T["title_font"], 26, bold=True, color=T["text_light"])


def _render_chart_for_ppt(slide_data: dict) -> bytes:
    """Render chart to PNG bytes for embedding in PPT."""
    stype = slide_data.get("type","bar_chart")
    categories = slide_data.get("categories",[])
    series_list = slide_data.get("series",[])
    fig, ax = plt.subplots(figsize=(7, 4), dpi=120, facecolor="#F4F6FB")
    ax.set_facecolor("#F4F6FB")

    if stype == "bar_chart":
        n = len(categories)
        x = np.arange(n)
        n_s = len(series_list)
        w = 0.7 / max(n_s, 1)
        for i, s in enumerate(series_list):
            vals = [float(v) if isinstance(v,(int,float)) else 0 for v in s.get("values",[])]
            col = _ACTIVE_THEME["chart_colors"][i % len(_ACTIVE_THEME["chart_colors"])]
            offset = (i - n_s/2 + 0.5) * w
            ax.bar(x + offset, vals[:n], width=w*0.9, color=col, label=s.get("name",""), alpha=0.85)
        ax.set_xticks(x); ax.set_xticklabels(categories, rotation=20, ha='right', fontsize=8)
    else:  # line
        n = len(categories)
        x = list(range(n))
        for i, s in enumerate(series_list):
            vals = [float(v) if isinstance(v,(int,float)) else 0 for v in s.get("values",[])]
            col = _ACTIVE_THEME["chart_colors"][i % len(_ACTIVE_THEME["chart_colors"])]
            ax.plot(x, vals[:n], marker='o', lw=2, color=col, label=s.get("name",""))
        ax.set_xticks(x); ax.set_xticklabels(categories, rotation=20, ha='right', fontsize=8)

    ax.set_title(slide_data.get("chart_title",""), fontsize=12, fontweight='bold')
    ax.set_xlabel(slide_data.get("x_label",""), fontsize=9)
    ax.set_ylabel(slide_data.get("y_label",""), fontsize=9)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    if len(series_list) > 1:
        ax.legend(fontsize=8)
    plt.tight_layout(pad=1.5)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _build_ppt_slide(prs, data: dict, slide_num: int):
    T = _ACTIVE_THEME
    stype = data.get("type","content")
    footer_y = _SLIDE_H - FOOTER_H

    if stype == "title":
        slide = _new_blank_slide(prs, T["bg_dark"])
        _add_textbox(slide, data.get("title",""), Inches(0.6), Inches(2.8),
                     Inches(11), Inches(1.2), T["title_font"], 40, bold=True, color=T["text_light"])
        sub = data.get("subtitle","")
        if sub:
            _add_textbox(slide, sub, Inches(0.6), Inches(4.1), Inches(10), Inches(0.6),
                         T["body_font"], 18, color=T["accent"])
        _add_footer(slide, slide_num)

    elif stype == "toc":
        slide = _new_blank_slide(prs, T["bg_light"])
        _add_header(slide, data.get("title","Agenda"))
        items = data.get("items", [])
        y = Inches(1.5)
        for i, item in enumerate(items):
            if y + Inches(0.55) > footer_y - Inches(0.1): break
            _add_textbox(slide, f"{i+1}.  {item}", Inches(0.6), y, Inches(11), Inches(0.45),
                         T["body_font"], 14, color=T["text_dark"])
            y += Inches(0.52)
        _add_footer(slide, slide_num)

    elif stype == "section":
        slide = _new_blank_slide(prs, T["bg_dark"])
        _add_textbox(slide, data.get("title",""), Inches(0.7), Inches(3.0),
                     Inches(11.5), Inches(1.2), T["title_font"], 36, bold=True, color=T["text_light"])
        _add_footer(slide, slide_num)

    elif stype == "content":
        slide = _new_blank_slide(prs, T["bg_light"])
        _add_header(slide, data.get("title",""))
        y = Inches(1.45)
        for bullet in data.get("bullets",[]):
            if y + Inches(0.44) > footer_y - Inches(0.1): break
            _add_textbox(slide, f"▸  {bullet.strip()}", Inches(0.6), y, Inches(11.8), Inches(0.44),
                         T["body_font"], 13, color=T["text_dark"])
            y += Inches(0.47)
        _add_footer(slide, slide_num)

    elif stype == "table":
        slide = _new_blank_slide(prs, T["bg_light"])
        _add_header(slide, data.get("title",""))
        headers = data.get("headers",[])
        rows = data.get("rows",[])
        if headers:
            all_rows = [headers] + rows
            cols = len(headers)
            col_w = Inches(12.6) / cols
            start_x = Inches(0.35)
            row_h = Inches(0.42)
            y = Inches(1.45)
            for r_idx, row in enumerate(all_rows):
                if y + row_h > footer_y - Inches(0.05): break
                for c_idx, cell in enumerate(row):
                    x = start_x + c_idx * col_w
                    is_hdr = r_idx == 0
                    bg = T["bg_dark"] if is_hdr else (PptxRGB(0xE8,0xEE,0xF8) if r_idx % 2 == 0 else T["bg_light"])
                    rect = slide.shapes.add_shape(1, x, y, col_w - Inches(0.02), row_h)
                    rect.fill.solid(); rect.fill.fore_color.rgb = bg
                    rect.line.color.rgb = PptxRGB(0xCC,0xCC,0xCC); rect.line.width = Pt(0.5)
                    _add_textbox(slide, str(cell), x + Inches(0.05), y + Inches(0.05),
                                 col_w - Inches(0.1), row_h - Inches(0.08),
                                 T["body_font"], 10 if is_hdr else 9,
                                 bold=is_hdr, color=T["text_light"] if is_hdr else T["text_dark"])
                y += row_h
        _add_footer(slide, slide_num)

    elif stype in ("bar_chart","line_chart"):
        slide = _new_blank_slide(prs, T["bg_light"])
        _add_header(slide, data.get("title",""))
        try:
            img_bytes = _render_chart_for_ppt(data)
            chart_stream = io.BytesIO(img_bytes)
            bullets = data.get("summary_bullets",[])
            chart_w = Inches(7.6) if bullets else Inches(12.0)
            slide.shapes.add_picture(chart_stream, Inches(0.35), Inches(1.4), chart_w, footer_y - Inches(1.4) - Inches(0.05))
            if bullets:
                y = Inches(2.0)
                for b in bullets:
                    if y + Inches(0.6) > footer_y - Inches(0.1): break
                    _add_textbox(slide, f"▸  {b.strip()}", Inches(8.1), y, Inches(4.8), Inches(0.58),
                                 T["body_font"], 11, color=T["text_dark"])
                    y += Inches(0.65)
        except Exception as e:
            logger.warning(f"Chart render failed: {e}")
        _add_footer(slide, slide_num)

    elif stype == "conclusion":
        slide = _new_blank_slide(prs, T["bg_dark"])
        _add_header(slide, data.get("title","Conclusion"))
        y = Inches(1.6)
        for b in data.get("bullets",[]):
            if y + Inches(0.5) > footer_y - Inches(0.1): break
            _add_textbox(slide, f"✦  {b.strip()}", Inches(0.6), y, Inches(11.5), Inches(0.5),
                         T["body_font"], 15, color=T["text_light"])
            y += Inches(0.58)
        _add_footer(slide, slide_num)

    elif stype == "thankyou":
        slide = _new_blank_slide(prs, T["bg_dark"])
        _add_textbox(slide, data.get("title","Thank You"), Inches(0.6), Inches(2.8),
                     Inches(12.0), Inches(1.4), T["title_font"], 52, bold=True,
                     color=T["text_light"], align=PP_ALIGN.CENTER)
        if data.get("message"):
            _add_textbox(slide, data["message"], Inches(0.6), Inches(4.2), Inches(12.0), Inches(0.8),
                         T["body_font"], 18, color=T["accent"], align=PP_ALIGN.CENTER)
        _add_footer(slide, slide_num)


def run_ppt_pipeline(combined_text: str, user_query: str, output_path: str):
    """Generate a combined analyzed PPTX from inspection text."""
    global _ACTIVE_THEME
    theme_key = random.choice(list(ALL_THEMES.keys()))
    _ACTIVE_THEME = ALL_THEMES[theme_key]
    logger.info(f"[ppt_pipeline] Using theme: {_ACTIVE_THEME['name']}")

    logger.info("[ppt_pipeline] Generating slide plan...")
    plan = _ppt_invoke(combined_text, user_query)

    prs = PPTXPresentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H

    title_data = plan.get("slides", [{}])[0] if plan.get("slides") else {}
    if title_data.get("type") == "title":
        title_data.setdefault("title",    plan.get("title","Inspection Analysis"))
        title_data.setdefault("subtitle", plan.get("subtitle",""))

    for i, slide_data in enumerate(plan.get("slides",[])):
        try:
            _build_ppt_slide(prs, slide_data, i + 1)
        except Exception as e:
            logger.warning(f"[ppt] Slide {i+1} failed: {e}")

    prs.save(output_path)
    logger.info(f"[ppt_pipeline] Done → {output_path}")
